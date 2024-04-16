import tkinter as tk
from tkinter import filedialog, Label
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageTk
import subprocess
import os
import sys


def open_pptx_file(ppt_file):
    try:
        if sys.platform == "win32":
            os.startfile(ppt_file)
        elif sys.platform == "darwin":
            subprocess.run(["open", ppt_file])
        elif sys.platform == "linux":
            subprocess.run(["xdg-open", ppt_file])
    except Exception as e:
        print(f"Failed to open file: {str(e)}")

def convert_pdf_to_ppt(pdf_file, ppt_file, status_callback):
    try:
        images = convert_from_path(pdf_file)
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # 16:9 aspect ratio width
        prs.slide_height = Inches(7.5)   # 16:9 aspect ratio height

        for image in images:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            img_path = 'temp_image.jpg'
            image.save(img_path, 'JPEG')
            
            # Add the picture to the slide, cover the entire slide
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(img_path)
        prs.save(ppt_file)
        status_callback("Conversion successful!", "success")
        open_pptx_file(ppt_file)
    except Exception as e:
        status_callback(f"An error occurred: {str(e)}", "error")

class PDFSelector(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title('PDF to PPT Converter')
        self.geometry('500x200')  # Set the window size

        # Add a button to open the file dialog
        self.open_file_btn = tk.Button(self, text='Open PDF File', command=self.open_file_dialog)
        self.open_file_btn.pack(pady=10)  # Add some vertical padding

        # Label to display the selected file path
        self.file_path_label = Label(self, text="", wraplength=300)
        self.file_path_label.pack(pady=10)  # Add some vertical padding

        # Add a button to start the conversion
        self.convert_btn = tk.Button(self, text="Convert to PPT", command=self.convert_to_ppt)
        self.convert_btn.pack(pady=10)

        # Status message label
        self.status_label = Label(self, text="", fg="black")
        self.status_label.pack(pady=10)

        self.file_path = None

    def open_file_dialog(self):
        # Open the file dialog to select a PDF
        file_path = filedialog.askopenfilename(filetypes=[('PDF files', '*.pdf')])
        if file_path:
            self.file_path = file_path
            self.file_path_label.config(text=file_path)  # Update the label with the file path

    def convert_to_ppt(self):
        if self.file_path:
            ppt_file = os.path.splitext(self.file_path)[0] + ".pptx"
            convert_pdf_to_ppt(self.file_path, ppt_file, self.update_status)
        else:
            self.update_status("Please select a PDF file first.")

    def update_status(self, message, status_type):
        self.status_label.config(text=message)
        if status_type == "success":
            self.status_label.config(fg="green")
        elif status_type == "error":
            self.status_label.config(fg="red")
        elif status_type == "warning":
            self.status_label.config(fg="orange")

# Create and run the application
app = PDFSelector()
app.mainloop()
